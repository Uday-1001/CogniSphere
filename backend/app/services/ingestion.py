from __future__ import annotations
import os
import logging
from typing import List, Optional, Dict, Callable
from langchain_community.document_loaders import PyMuPDFLoader, TextLoader
from langchain_core.documents import Document
from langchain_text_splitters import RecursiveCharacterTextSplitter, MarkdownHeaderTextSplitter
from .metadata import MetadataEnricher
from ..config.settings import settings

logger = logging.getLogger(__name__)


try:
    from .ocr_pipeline import PDFOCRPipeline
    ocr_pipeline = PDFOCRPipeline(
        dpi=settings.OCR_DPI,
        languages=settings.OCR_LANGUAGE.split(","),
        char_threshold=settings.OCR_SCANNED_CHAR_THRESHOLD,
        max_workers=settings.OCR_MAX_WORKERS or 4,
    )
    ocr_available = True
except ImportError as _ocr_import_err:
    ocr_available = False
    ocr_pipeline = None
    logger.warning(
        "OCR pipeline unavailable (%s). Scanned PDFs will not be processed.",
        _ocr_import_err,
    )


try:
    from .docling_parser import parse_with_docling
    docling_available = True
except ImportError as docling_import_err:
    docling_available = False
    parse_with_docling = None
    logger.warning(
        "Docling unavailable (%s). Falling back to PyMuPDFLoader for PDFs.",
        docling_import_err,
    )


Markdown_headers = [
    ("#",   "heading_1"),
    ("##",  "heading_2"),
    ("###", "heading_3"),
]


def make_pymupdf_documents(file_path: str) -> List[Document]:
    loader = PyMuPDFLoader(file_path)
    docs = list(loader.lazy_load())
    for doc in docs:
        raw_page = doc.metadata.pop("page", None)
        doc.metadata.setdefault("page_number", (raw_page + 1) if raw_page is not None else None)
        doc.metadata.setdefault("section",       None)
        doc.metadata["parser_used"]   = "pymupdf"
        doc.metadata["is_ocr"]        = False
        doc.metadata["document_type"] = "pdf_digital"
    return docs


class IngestionService:

    enriched_metadata = MetadataEnricher()

    supported_extensions = {
        "video":    [".mp4", ".mov", ".mkv", ".avi", ".webm"],
        "audio":    [".mp3", ".wav", ".m4a", ".flac"],
        "document": [".pdf", ".docx", ".pptx", ".txt"],
    }


    def get_file_type(self, filename: str) -> Optional[str]:
        ext = os.path.splitext(filename)[1].lower()
        for file_type, extensions in self.supported_extensions.items():
            if ext in extensions:
                return file_type
        return None

    def load_document(
        self,
        file_path: str,
        file_type: str,
        progress_callback: Optional[Callable] = None,
    ) -> List[Document]:
        ext = os.path.splitext(file_path)[1].lower()

        if ext == ".pdf":
            return self.load_pdf(file_path, progress_callback)
        elif ext == ".docx":
            return self.load_docx(file_path)
        elif ext == ".pptx":
            return self.load_pptx(file_path)
        elif ext == ".txt":
            return self.load_txt(file_path)
        elif file_type in ["video", "audio"]:
            return []
        else:
            raise ValueError(f"Unsupported file extension: {ext}")

    def load_transcript(self, transcript_path: Optional[str]) -> List[Document]:
        if not transcript_path:
            return []
        return TextLoader(transcript_path, encoding="utf-8").load()

    def split_documents(
        self,
        documents: List[Document],
        chunk_size: int = 1500,
        chunk_overlap: int = 300,
    ) -> List[Document]:
        recursive_splitter = RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            add_start_index=True,
        )
        md_splitter = MarkdownHeaderTextSplitter(
            headers_to_split_on=Markdown_headers,
            strip_headers=False,
        )

        result: List[Document] = []

        for doc in documents:
            if doc.metadata.get("parser_used") == "docling":
                try:
                    header_chunks = md_splitter.split_text(doc.page_content)
                except Exception as md_exc:
                    logger.warning(
                        "MarkdownHeaderTextSplitter failed, falling back to "
                        "RecursiveCharacterTextSplitter: %s", md_exc
                    )
                    header_chunks = [doc]

                for chunk in header_chunks:
                    chunk.metadata = {**doc.metadata, **chunk.metadata}

                further_split = recursive_splitter.split_documents(header_chunks)
                result.extend(further_split)
            else:
                result.extend(recursive_splitter.split_documents([doc]))

        return result

    def process_document(
        self,
        file_path: str,
        file_id: int,
        filename: str,
        file_type: str,
        transcript_path: Optional[str] = None,
        segment_timestamps: Optional[List[Dict]] = None,
        progress_callback: Optional[Callable] = None,
    ) -> List[Document]:

        if file_type in ["video", "audio"]:
            documents = self.load_transcript(transcript_path)
            for doc in documents:
                doc.metadata["source_type"] = file_type
                doc.metadata["filename"]    = filename
        else:
            documents = self.load_document(file_path, file_type, progress_callback)
            for doc in documents:
                doc.metadata["source_type"] = file_type
                doc.metadata["filename"]    = filename

        split_docs = self.split_documents(documents)

        enriched: List[Document] = []
        for i, doc in enumerate(split_docs):
            if file_type in ["video", "audio"] and segment_timestamps:
                ts_start, ts_end = self.enriched_metadata.extract_timestamps_from_segments(
                    doc.page_content, segment_timestamps
                )
                enriched_doc = self.enriched_metadata.enrich(
                    doc, file_id, filename, i,
                    timestamp_start=ts_start,
                    timestamp_end=ts_end,
                    parser_used=doc.metadata.get("parser_used"),
                    is_ocr=doc.metadata.get("is_ocr"),
                    document_type=doc.metadata.get("document_type"),
                    ocr_confidence=doc.metadata.get("ocr_confidence"),
                    section=doc.metadata.get("section"),
                )
            else:
                enriched_doc = self.enriched_metadata.enrich(
                    doc, file_id, filename, i,
                    page_number=doc.metadata.get("page_number"),
                    parser_used=doc.metadata.get("parser_used"),
                    is_ocr=doc.metadata.get("is_ocr"),
                    document_type=doc.metadata.get("document_type"),
                    ocr_confidence=doc.metadata.get("ocr_confidence"),
                    section=doc.metadata.get("section"),
                )
            enriched.append(enriched_doc)

        return enriched


    def load_pdf(
        self,
        file_path: str,
        progress_callback: Optional[Callable] = None,
    ) -> List[Document]:
        filename = os.path.basename(file_path)

        if ocr_available and ocr_pipeline is not None:
            try:
                if ocr_pipeline.is_scanned(file_path):
                    logger.info("'%s' is scanned — routing to PaddleOCR.", filename)
                    return self.run_ocr(file_path, progress_callback)
            except Exception as scan_exc:
                logger.warning(
                    "Scan-detection failed for '%s' (%s) — assuming digital.",
                    filename, scan_exc,
                )

        if docling_available and parse_with_docling is not None:
            try:
                docs = parse_with_docling(file_path)
                total_text = " ".join(d.page_content for d in docs).strip()
                if len(total_text) < 100:
                    raise RuntimeError(
                        f"Docling extracted only {len(total_text)} characters — "
                        "content appears empty; trying fallback."
                    )
                return docs
            except Exception as docling_exc:
                logger.warning(
                    "Docling failed for '%s' (%s) — falling back to PyMuPDFLoader.",
                    filename, docling_exc,
                )

        logger.info("'%s' — using PyMuPDFLoader (fallback).", filename)
        docs = make_pymupdf_documents(file_path)

        total_text = " ".join(d.page_content for d in docs).strip()
        if len(total_text) < 100 and ocr_available and ocr_pipeline is not None:
            logger.info(
                "'%s' yielded only %d chars via PyMuPDF — attempting OCR fallback.",
                filename, len(total_text),
            )
            ocr_docs = self.run_ocr(file_path, progress_callback)
            if ocr_docs:
                return ocr_docs

        return docs


    def load_docx(self, file_path: str) -> List[Document]:
        filename = os.path.basename(file_path)

        if docling_available and parse_with_docling is not None:
            try:
                return parse_with_docling(file_path)
            except Exception as exc:
                logger.warning(
                    "Docling failed for DOCX '%s' (%s) — falling back to python-docx parser.",
                    filename, exc,
                )

        try:
            import docx

            doc = docx.Document(file_path)
            full_text: List[str] = []
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    full_text.append(text)
            for table in doc.tables:
                for row in table.rows:
                    row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_text:
                        full_text.append(" | ".join(row_text))

            content = "\n\n".join(full_text)
            return [
                Document(
                    page_content=content,
                    metadata={
                        "page_number": 1,
                        "section": None,
                        "parser_used": "python-docx",
                        "is_ocr": False,
                        "document_type": "docx",
                    },
                )
            ]
        except Exception as exc:
            logger.error("Failed to parse DOCX '%s': %s", filename, exc)
            return [
                Document(
                    page_content="",
                    metadata={
                        "page_number": 1,
                        "section": None,
                        "parser_used": "docx_fallback",
                        "is_ocr": False,
                        "document_type": "docx",
                    },
                )
            ]


    def load_pptx(self, file_path: str) -> List[Document]:
        filename = os.path.basename(file_path)

        if docling_available and parse_with_docling is not None:
            try:
                return parse_with_docling(file_path)
            except Exception as exc:
                logger.warning(
                    "Docling failed for PPTX '%s' (%s) — falling back to python-pptx parser.",
                    filename, exc,
                )

        try:
            import pptx

            prs = pptx.Presentation(file_path)
            slide_docs: List[Document] = []
            for slide_num, slide in enumerate(prs.slides, start=1):
                slide_text: List[str] = []
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text_frame = getattr(shape, "text_frame", None)
                        if text_frame:
                            for paragraph in text_frame.paragraphs:
                                text = paragraph.text.strip()
                                if text:
                                    slide_text.append(text)
                if slide_text:
                    slide_docs.append(
                        Document(
                            page_content="\n".join(slide_text),
                            metadata={
                                "page_number": slide_num,
                                "section": None,
                                "parser_used": "python-pptx",
                                "is_ocr": False,
                                "document_type": "pptx",
                            },
                        )
                    )
            if not slide_docs:
                slide_docs = [
                    Document(
                        page_content="",
                        metadata={
                            "page_number": 1,
                            "section": None,
                            "parser_used": "python-pptx",
                            "is_ocr": False,
                            "document_type": "pptx",
                        },
                    )
                ]
            return slide_docs
        except Exception as exc:
            logger.error("Failed to parse PPTX '%s': %s", filename, exc)
            return [
                Document(
                    page_content="[PPTX parsing failed]",
                    metadata={
                        "page_number": 1,
                        "section": None,
                        "parser_used": "pptx_fallback",
                        "is_ocr": False,
                        "document_type": "pptx",
                    },
                )
            ]


    def load_txt(self, file_path: str) -> List[Document]:
        docs = TextLoader(file_path, encoding="utf-8").load()
        for doc in docs:
            doc.metadata.setdefault("page_number",   1)
            doc.metadata.setdefault("section",       None)
            doc.metadata["parser_used"]   = "textloader"
            doc.metadata["is_ocr"]        = False
            doc.metadata["document_type"] = "txt"
        return docs


    def run_ocr(
        self,
        file_path: str,
        progress_callback: Optional[Callable] = None,
    ) -> List[Document]:
        if not ocr_available or ocr_pipeline is None:
            raise RuntimeError(
                "OCR is required for this file but PaddleOCR / PyMuPDF are not installed. "
                "Run: pip install paddlepaddle paddleocr PyMuPDF"
            )
        filename = os.path.basename(file_path)
        return ocr_pipeline.process(file_path, filename, progress_callback)


ingestion_service = IngestionService()